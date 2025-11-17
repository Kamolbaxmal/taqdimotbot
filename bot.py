import os
import io
import logging
import wikipedia
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt
from telegram import Update, InputFile
from telegram.ext import ApplicationBuilder, CommandHandler, ContextTypes

# ========== WIKI TILI ==========
wikipedia.set_lang("uz")

# ========== ENV ================
load_dotenv()
BOT_TOKEN = os.getenv("BOT_TOKEN")

logging.basicConfig(level=logging.INFO)

# =================================
# PPT YARATISH FUNKSIYASI
# =================================

def create_presentation(topic: str) -> bytes:
    prs = Presentation()

    try:
        summary = wikipedia.summary(topic, sentences=30)
    except:
        summary = "Bu mavzu bo‘yicha ma’lumot topilmadi."

    # Matnni 10 bo‘lakka bo‘lamiz
    parts = []
    words = summary.split()
    chunk_size = len(words) // 10 + 1

    for i in range(0, len(words), chunk_size):
        parts.append(" ".join(words[i:i + chunk_size]))

    for idx, text in enumerate(parts[:10]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"{topic} — {idx+1}-slayd"
        p = body.text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(20)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# =================================
# /start
# =================================

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Salom!\n"
        "/make <mavzu> — mavzu bo‘yicha 10 slaydli taqdimot tayyorlayman.\n"
        "Masalan:\n"
        "/make Sun'iy intellekt"
    )


# =================================
# /make
# =================================

async def make_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not context.args:
        await update.message.reply_text("❗ Foydalanish: /make <mavzu>")
        return

    topic = " ".join(context.args)
    await update.message.reply_text(f"⏳ '{topic}' bo‘yicha slaydlar yaratilmoqda...")

    ppt_bytes = create_presentation(topic)

    file = io.BytesIO(ppt_bytes)
    file.name = f"{topic}.pptx"

    await update.message.reply_document(InputFile(file))
    await update.message.reply_text("✅ Tayyor!")


# =================================
# MAIN
# =================================

def main():
    app = ApplicationBuilder().token(BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("make", make_presentation))

    app.run_polling()


if __name__ == "__main__":
    print("Bot ishga tushdi...")
    main()
